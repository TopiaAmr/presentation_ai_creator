from pptx import Presentation
from pptx.util import Inches, Pt  # Inches for size/position, Pt for font size

# Import other necessary parts like shapes, charts if needed

# --- Slide Layout Mapping (adjust based on your default template) ---
# Run this code once locally to find the layout indices for your template:
# prs_check = Presentation() # Use your template if desired: Presentation('template.pptx')
# for i, layout in enumerate(prs_check.slide_layouts):
#     print(f"Index: {i}, Name: {layout.name}")
# Common indices often are: 0=Title, 1=Title+Content, 5=Blank, etc.
TITLE_SLIDE_LAYOUT = 0
TITLE_AND_CONTENT_LAYOUT = 1
BLANK_LAYOUT = 5
# --- End Layout Mapping ---


def create_presentation_from_data(data):
    """
    Generates a PowerPoint presentation object from structured data.
    Input 'data' structure (example):
    {
        "title": "Optional Presentation Title (Not used on slides directly yet)",
        "slides": [
            {
                "layout_type": "title_slide", # Or 'title_content', 'blank' etc.
                "title": "Slide 1 Title",
                "subtitle": "Optional Subtitle for Title Slide"
            },
            {
                "layout_type": "title_content",
                "title": "Slide 2: Key Points",
                "content": [ # List of strings for bullet points
                    "First point.",
                    "Second point.",
                    "  Indented sub-point." # Indentation detected by leading spaces
                ]
            },
            {
                "layout_type": "title_content",
                 "title": "Slide 3: Paragraph",
                 "content": "This is a single paragraph of text." # Can also be a single string
            }
            # Add more slide types/data as needed (images, charts etc.)
        ]
    }
    """
    prs = (
        Presentation()
    )  # Consider using a template: Presentation('your_template.pptx')

    for slide_data in data.get("slides", []):
        layout_type = slide_data.get("layout_type", "title_content").lower()
        slide_layout = None

        # Select the layout based on input type
        if layout_type == "title_slide":
            slide_layout = prs.slide_layouts[TITLE_SLIDE_LAYOUT]
        elif layout_type == "title_content":
            slide_layout = prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT]
        elif layout_type == "blank":
            slide_layout = prs.slide_layouts[BLANK_LAYOUT]
        else:  # Default to title and content if unknown
            print(f"Warning: Unknown layout type '{layout_type}', defaulting.")
            slide_layout = prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT]

        slide = prs.slides.add_slide(slide_layout)

        # --- Populate Placeholders ---
        # Check if title placeholder exists using a try/except block instead of has_title attribute
        title_placeholder = None
        try:
            title_placeholder = slide.shapes.title
        except AttributeError:
            title_placeholder = None

        # Placeholder indices can vary - check your template or guess common ones
        content_placeholder = None
        subtitle_placeholder = None
        try:
            if layout_type == "title_slide":
                subtitle_placeholder = slide.placeholders[
                    1
                ]  # Often index 1 for subtitle
            elif layout_type == "title_content":
                content_placeholder = slide.placeholders[
                    1
                ]  # Often index 1 for main content
        except KeyError:
            print(
                f"Warning: Could not find expected placeholder on slide layout '{layout_type}'"
            )

        # Add Title
        if title_placeholder and slide_data.get("title"):
            title_placeholder.text = slide_data["title"]

        # Add Subtitle (for Title Slide)
        if subtitle_placeholder and slide_data.get("subtitle"):
            subtitle_placeholder.text = slide_data["subtitle"]

        # Add Content (Bullets or Paragraph)
        if content_placeholder and slide_data.get("content"):
            content = slide_data["content"]
            tf = content_placeholder.text_frame
            tf.clear()  # Clear existing placeholder text

            if isinstance(content, list):  # Handle as bullet points
                for i, point in enumerate(content):
                    level = 0
                    text = point
                    # Simple indentation check
                    if point.startswith("    "):  # 4 spaces = level 2
                        level = 2
                        text = point.lstrip()
                    elif point.startswith("  "):  # 2 spaces = level 1
                        level = 1
                        text = point.lstrip()

                    if i == 0:  # First point uses existing paragraph
                        p = tf.paragraphs[0]
                        p.text = text
                        p.level = level
                    else:  # Add new paragraphs
                        p = tf.add_paragraph()
                        p.text = text
                        p.level = level
            elif isinstance(content, str):  # Handle as single paragraph
                tf.paragraphs[0].text = content
            else:
                print(f"Warning: Unsupported content type: {type(content)}")

    return prs  # Return the Presentation object


# --- Example Usage (for testing locally) ---
if __name__ == "__main__":
    sample_data = {
        "slides": [
            {
                "layout_type": "title_slide",
                "title": "AI Presentation Demo",
                "subtitle": "Generated via API",
            },
            {
                "layout_type": "title_content",
                "title": "How it Works",
                "content": [
                    "N8N generates content JSON.",
                    "API receives JSON.",
                    "Python script uses python-pptx.",
                    "  Generates .pptx file.",
                    "API returns the file.",
                ],
            },
            {
                "layout_type": "title_content",
                "title": "Benefits",
                "content": "Automation, Consistency, Speed.",
            },
        ]
    }
    presentation = create_presentation_from_data(sample_data)
    output_filename = "generated_presentation.pptx"
    presentation.save(output_filename)
    print(f"Presentation saved as {output_filename}")
